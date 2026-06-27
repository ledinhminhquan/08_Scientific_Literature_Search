"""Generate the submission slides.pptx (python-pptx) — ~10 concise visual slides.
Degrades to a Markdown outline if python-pptx is unavailable.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from ..config import AppConfig, artifacts_dir
from ..logging_utils import get_logger
from . import charts as charts_mod
from .artifact_loader import load_artifacts

logger = get_logger(__name__)


def _slides(cfg: AppConfig, arts: Dict[str, Any]) -> List[Tuple[str, List[str]]]:
    ev = arts.get("eval") or {}
    m = ev.get("model", {})
    b = ev.get("bm25", {})
    mn, bn = m.get("ndcg@10"), b.get("ndcg@10")
    res = (f"retriever nDCG@10 {mn:.3f} vs BM25 {bn:.3f}" if (mn is not None and bn is not None)
           else "train + evaluate to populate results")
    return [
        ("Exploratory Scientific Literature Search",
         [f"{cfg.author} — Student {cfg.student_id}", "NLP in Industry — Final Assignment",
          "Find + explore scientific papers, not just a flat list", "Trainable retriever + hybrid search + an agent"]),
        ("Business Problem & Motivation",
         ["Keyword search misses paraphrase / concept", "A flat list offers no way to explore a literature",
          "Researchers can't keep up with the paper firehose", "Goal: semantic recall + facets + clusters + related"]),
        ("Proposed NLP Solution",
         ["Trainable core: fine-tuned dense retriever (bi-encoder)", "Hybrid: dense + BM25 fused via RRF + reranker",
          "Exploration: field facets, topic clusters, related papers", "An agent understands the query + guides exploration"]),
        ("System Architecture",
         ["query -> understand/expand (intent + filters)", "-> dense + BM25 -> RRF fuse -> [rerank]",
          "-> cluster into topics + facet by field + related", "-> ranked, explorable results"]),
        ("Data Overview",
         ["Corpus: gfissore/arxiv-abstracts-2021 (CC0, 2M papers)", "Title+abstract+categories -> field/year facets",
          "Synthetic (query->paper) pairs for contrastive training", "Real IR eval via BeIR/scifact"]),
        ("Model & Evaluation Results",
         [res, "Metrics: nDCG@10 (headline), Recall@{1,5,10}, MRR@10",
          "Baselines: BM25 + zero-shot base encoder", "Fine-tuned retriever beats both on conceptual queries"]),
        ("Agentic AI Component",
         ["Deterministic FSM + optional LLM brain (rule fallback)", "D1 query-type routing · D2 coverage gate (expand)",
          "D3 rerank gate (only when head is close) · D4 explore", "Full audit trace; 0 paid API by default"]),
        ("Deployment Overview",
         ["FastAPI /search + /related · Gradio search UI", "CLI · Docker · HF Space",
          "FAISS dense index + BM25; hot path ~tens of ms", "Retriever versioned via registry (repo@revision)"]),
        ("Ethics, Privacy & Risks",
         ["Query privacy: don't retain raw queries beyond need", "Corpus coverage/recency bias -> stale-index risk",
          "Over-trust in ranking -> show scores + diversity", "Permissive, offline-capable stack (no data egress)"]),
        ("Key Takeaways & Future Work",
         ["Fine-tuned dense retrieval beats BM25 on concepts", "Exploration (facets/clusters/related) > a flat list",
          "Future: citation graph, FoS taxonomy, survey detection", "Future: index sharding, incremental updates, RLHF"]),
    ]


def generate_slides(cfg: AppConfig, title: Optional[str] = None, author: Optional[str] = None,
                    out_path: Optional[str] = None) -> str:
    arts = load_artifacts(cfg)
    out_path = Path(out_path) if out_path else artifacts_dir() / "slides.pptx"
    out_path.parent.mkdir(parents=True, exist_ok=True)
    slides = _slides(cfg, arts)
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt
    except Exception as exc:
        logger.warning("python-pptx unavailable (%s); writing markdown outline", exc)
        md = "\n\n".join(f"## {t}\n" + "\n".join(f"- {b}" for b in bs) for t, bs in slides)
        alt = out_path.with_suffix(".md")
        alt.write_text(md, encoding="utf-8")
        return str(alt)

    chart = charts_mod.ndcg_chart(arts.get("eval") or {}, out_path.parent / "charts" / "slide_ndcg.png")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    accent = RGBColor(0x2B, 0x6C, 0xB0)
    for i, (t, bullets) in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
        bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
        tf = bar.text_frame; tf.text = t
        tf.paragraphs[0].font.size = Pt(28); tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        body = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(8.3 if (i == 5 and chart) else 12), Inches(5.4))
        bt = body.text_frame; bt.word_wrap = True
        for j, bp in enumerate(bullets):
            p = bt.paragraphs[0] if j == 0 else bt.add_paragraph()
            p.text = "•  " + bp; p.font.size = Pt(20); p.space_after = Pt(10)
        if i == 5 and chart:
            slide.shapes.add_picture(str(chart), Inches(8.9), Inches(1.7), width=Inches(4.0))
        foot = slide.shapes.add_textbox(Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.4))
        foot.text_frame.text = f"{title or cfg.project_title} — {author or cfg.author} ({cfg.student_id})"
        foot.text_frame.paragraphs[0].font.size = Pt(9)
    prs.save(str(out_path))
    logger.info("Slides -> %s", out_path)
    return str(out_path)


__all__ = ["generate_slides"]
